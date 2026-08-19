"""Build the deck as a PowerPoint file.

    pip install python-pptx
    python deck/build_pptx.py

Writes deck/verification-budget.pptx.

The HTML deck in index.html is the primary artifact, and it is the one to
present from if you can. This exists because conference AV desks ask for a pptx,
and because some people want to edit the slides in the tool they already have.

Layout is written out here rather than templated, because the visual direction is
the point and a stock template would undo it. Speaker notes are read from the
same `deck/slides/*.md` files the HTML deck uses, so the two cannot drift apart.

Visual direction: derived from the information design principles of Edward Tufte.
Content first, high content-to-ink ratio, hairline rules instead of boxes, one
accent at a time, serif for prose and a consistent treatment for every number.
Dark title and closing, light content in between. No template furniture, no
decorative stripes, no logo on every slide.
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
ROOT = HERE.parent
SLIDES = HERE / "slides"
ASSETS = HERE / "assets"
OUTPUT = HERE / "verification-budget.pptx"

sys.path.insert(0, str(ROOT))

# ---------------------------------------------------------------------------
# Meter, the repository palette
# ---------------------------------------------------------------------------

INK = RGBColor(0x10, 0x14, 0x18)
SLATE = RGBColor(0x2A, 0x33, 0x3B)
STEEL = RGBColor(0x4A, 0x55, 0x5F)
MIST = RGBColor(0x6B, 0x77, 0x80)
LINE = RGBColor(0xD8, 0xD4, 0xCB)
SHELL = RGBColor(0xF2, 0xF1, 0xEE)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
BRAND = RGBColor(0x1F, 0x3A, 0x34)

CLASS_COLOUR = {
    "A": RGBColor(0x2E, 0x7D, 0x6B),
    "B": RGBColor(0x3C, 0x6E, 0x9F),
    "C": RGBColor(0xB8, 0x84, 0x3A),
    "D": RGBColor(0xA3, 0x44, 0x32),
}
OVER = CLASS_COLOUR["D"]
OK = RGBColor(0x1F, 0x7A, 0x5C)

ON_DARK = RGBColor(0xF7, 0xF6, 0xF3)
ON_DARK_MUTED = RGBColor(0xC9, 0xC6, 0xBF)
ON_DARK_ACCENT = RGBColor(0x7F, 0xD8, 0xC2)
DARK_LINE = RGBColor(0x2B, 0x32, 0x38)

# Both ship with Office and render true to width, so a fit that looks right
# here is right in the room.
SERIF = "Cambria"
SANS = "Calibri"

W = 13.333
H = 7.5
MARGIN = 0.85
CONTENT_W = W - MARGIN * 2


# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------


def blank(prs: Presentation, background: RGBColor = PAPER):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = background
    return slide


def text(
    slide,
    body: str,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int = 15,
    font: str = SANS,
    colour: RGBColor = SLATE,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing: float = 1.25,
    space_after: int = 0,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0

    for index, chunk in enumerate(body.split("\n")):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        para.line_spacing = line_spacing
        para.space_after = Pt(space_after)
        run = para.add_run()
        run.text = chunk
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = colour
    return box


def rich(
    slide,
    parts: list[tuple[str, dict]],
    left: float,
    top: float,
    width: float,
    height: float,
    size: int = 15,
    font: str = SANS,
    colour: RGBColor = SLATE,
    line_spacing: float = 1.25,
):
    """One paragraph made of differently styled runs."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    para = frame.paragraphs[0]
    para.line_spacing = line_spacing
    for content, style in parts:
        run = para.add_run()
        run.text = content
        run.font.name = style.get("font", font)
        run.font.size = Pt(style.get("size", size))
        run.font.bold = style.get("bold", False)
        run.font.italic = style.get("italic", False)
        run.font.color.rgb = style.get("colour", colour)
    return box


def hairline(slide, left: float, top: float, width: float, colour: RGBColor = LINE, weight: float = 0.9):
    """A rule. Used to separate rows of a table, which is what a rule is for."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Emu(int(weight * 9525))
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def card(slide, left: float, top: float, width: float, height: float,
         fill: RGBColor = SHELL, radius: float = 0.04):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    shape.adjustments[0] = radius
    return shape


def chip(slide, letter: str, left: float, top: float, size: float = 0.32):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CLASS_COLOUR[letter]
    shape.line.fill.background()
    shape.shadow.inherit = False
    shape.adjustments[0] = 0.22
    frame = shape.text_frame
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = letter
    run.font.name = SANS
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = ON_DARK
    return shape


def eyebrow(slide, label: str):
    text(slide, label.upper(), MARGIN, 0.52, CONTENT_W, 0.28,
         size=10, font=SANS, colour=MIST, bold=True)


def title(slide, heading: str, top: float = 1.02, size: int = 33, width: float = 10.4,
          colour: RGBColor = INK):
    # 1.24in fits two lines at 33pt. Content on every slide starts at 2.40 or
    # later, so the box cannot run into it.
    return text(slide, heading, MARGIN, top, width, 1.24,
                size=size, font=SERIF, colour=colour, bold=True, line_spacing=1.06)


HOST_LOGO_CANDIDATES = ("gsdc-logo.png", "gsdc-logo.jpg", "gsdc-logo.jpeg")
HOST_NAME = "GSDC Certified Learning Masterclass Series"


def host_logo_path():
    for name in HOST_LOGO_CANDIDATES:
        path = ASSETS / name
        if path.exists():
            return path
    return None


def host_logo(slide, top: float = 0.34, height: float = 0.40) -> None:
    """Top right, on light slides only. Height fixed, width follows the aspect."""
    path = host_logo_path()
    if path is None:
        return
    from PIL import Image  # optional; fall back to a fixed box without it

    try:
        with Image.open(path) as im:
            ratio = im.width / im.height
    except Exception:
        ratio = 3.2
    width = height * ratio
    slide.shapes.add_picture(str(path), Inches(W - MARGIN - width),
                             Inches(top), Inches(width), Inches(height))


def footer(slide, number: int, total: int):
    hairline(slide, MARGIN, H - 0.72, CONTENT_W)
    text(slide, "github.com/hotragn/verb     x.com/hotragn     linkedin.com/in/hotragn-pettugani",
         MARGIN, H - 0.60, 8.4, 0.26, size=9, font=SANS, colour=MIST)
    text(slide, f"{number} / {total}", W - MARGIN - 1.2, H - 0.60, 1.2, 0.26,
         size=9, font=SANS, colour=MIST, align=PP_ALIGN.RIGHT)


def part_divider(slide, part: str, heading: str, lead: str) -> None:
    """Dark full-bleed divider. No footer, like the title and the closing."""
    text(slide, part.upper(), MARGIN, 2.30, 3.0, 0.3, size=12, font=SANS,
         colour=ON_DARK_ACCENT, bold=True)
    hairline(slide, MARGIN, 2.72, 2.0, DARK_LINE, 1.2)
    text(slide, heading, MARGIN, 3.02, 9.5, 1.1, size=44, font=SERIF,
         colour=ON_DARK, bold=True, line_spacing=1.06)
    text(slide, lead, MARGIN, 4.34, 7.6, 1.1, size=18, font=SERIF,
         colour=ON_DARK_MUTED, line_spacing=1.32)


def table(slide, top: float, cols: list[tuple[float, float, str]],
          rows: list[list[str]], pitch: float = 0.56, size: float = 13.5,
          bold_first: bool = False, row_colours: list = None):
    """Hairline table. Header underlined heavily, rows lightly, nothing boxed."""
    for left, width, head in cols:
        text(slide, head.upper(), left, top, width, 0.26, size=9.5, font=SANS,
             colour=MIST, bold=True)
    hairline(slide, cols[0][0], top + 0.28, sum(w for _, w, _ in cols)
             + (cols[-1][0] - cols[0][0] - sum(w for _, w, _ in cols[:-1])), INK, 1.4)
    y = top + 0.46
    for index, row in enumerate(rows):
        colour = (row_colours[index] if row_colours else None) or SLATE
        for (left, width, _), cell in zip(cols, row):
            first = cell is row[0]
            text(slide, cell, left, y, width, pitch - 0.08, size=size, font=SANS,
                 colour=INK if (first and bold_first) else colour,
                 bold=first and bold_first, line_spacing=1.16)
        y += pitch
        hairline(slide, cols[0][0], y - 0.10,
                 cols[-1][0] + cols[-1][1] - cols[0][0])
    return y


def stat(slide, left: float, top: float, width: float, label: str, value: str,
         note: str, value_colour: RGBColor = INK, value_size: int = 46):
    hairline(slide, left, top, width, INK, 1.4)
    text(slide, label.upper(), left, top + 0.12, width, 0.24, size=9.5, font=SANS,
         colour=MIST, bold=True)
    text(slide, value, left, top + 0.38, width, 0.78, size=value_size, font=SERIF,
         colour=value_colour, bold=True, line_spacing=1.0)
    text(slide, note, left, top + 1.16, width, 0.9, size=11.5, font=SANS, colour=MIST,
         line_spacing=1.24)


def numbered(slide, left: float, top: float, width: float, index: str,
             head: str, body: str, head_size: int = 15, gap: float = 0.42):
    text(slide, index, left, top + 0.02, gap, 0.3, size=12, font=SANS, colour=MIST, bold=True)
    rich(slide, [(head, {"bold": True, "colour": INK, "size": head_size}),
                 ("  " + body, {"colour": SLATE, "size": head_size - 0.5})],
         left + gap, top, width - gap, 0.9, line_spacing=1.26)


# ---------------------------------------------------------------------------
# Speaker notes, read from the same sources the HTML deck uses
# ---------------------------------------------------------------------------


def load_notes() -> list[str]:
    notes: list[str] = []
    for path in sorted(SLIDES.glob("*.md")):
        body = path.read_text(encoding="utf-8")
        chunk = body.split("\n@notes\n", 1)[1] if "\n@notes\n" in body else ""
        cleaned = re.sub(r"\*\*(.+?)\*\*", r"\1", chunk.strip())
        cleaned = re.sub(r"\*(.+?)\*", r"\1", cleaned)
        notes.append(cleaned)
    return notes


def attach_notes(slide, body: str) -> None:
    if body:
        slide.notes_slide.notes_text_frame.text = body


# ---------------------------------------------------------------------------
# The slides
# ---------------------------------------------------------------------------

TOTAL = 27


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    notes = load_notes()
    n = 0

    def new(background: RGBColor = PAPER, section: str | None = None, chrome: bool = True):
        nonlocal n
        n += 1
        slide = blank(prs, background)
        if section:
            eyebrow(slide, section)
            hairline(slide, MARGIN, 0.86, CONTENT_W)
        if chrome:
            footer(slide, n, TOTAL)
            host_logo(slide)
        attach_notes(slide, notes[n - 1] if n - 1 < len(notes) else "")
        return slide

    # 1. Title -----------------------------------------------------------
    s = new(INK, chrome=False)
    text(s, HOST_NAME.upper(), MARGIN, 1.42, 8.0, 0.3, size=11, font=SANS,
         colour=ON_DARK_ACCENT, bold=True)
    hairline(s, MARGIN, 2.05, 1.0, ON_DARK, 2.6)
    text(s, "The Verification Budget", MARGIN, 2.42, 10.0, 2.0,
         size=54, font=SERIF, colour=ON_DARK, bold=True, line_spacing=1.05)
    text(s, "An operating model for autonomous AI in project delivery\nand end to end PMO.",
         MARGIN, 4.36, 9.0, 1.0, size=19, font=SERIF, colour=ON_DARK_MUTED, line_spacing=1.3)
    hairline(s, MARGIN, 5.72, CONTENT_W, DARK_LINE, 1.2)
    text(s, "Hotragn Pettugani", MARGIN, 5.92, 3.4, 0.3, size=13, font=SANS, colour=ON_DARK)
    text(s, "github.com/hotragn/verb          x.com/hotragn          linkedin.com/in/hotragn-pettugani",
         MARGIN + 3.5, 5.92, 8.5, 0.3, size=13, font=SANS, colour=ON_DARK_ACCENT)

    # 2. Part 1 divider ---------------------------------------------------
    s = new(INK, chrome=False)
    part_divider(s, "Part 1", "Why the constraint moved",
                 "Production stopped being scarce. Review did not.")

    # 3. Thirty years ----------------------------------------------------
    s = new(section="The constraint")
    title(s, "For thirty years, the scarce thing\nwas getting the work done.")
    text(s, "Somebody had to write the status report, chase the risk owners, rebuild the "
            "forecast after the change request landed. Capacity meant counting the people "
            "who could make those things.",
         MARGIN, 2.72, 7.1, 1.3, size=16, line_spacing=1.42)
    rich(s, [("That is over. ", {"bold": True, "colour": INK, "size": 20}),
             ("An agent produces a defensible status pack for forty projects in the time "
              "it takes you to open the file.", {"size": 20, "colour": SLATE})],
         MARGIN, 4.32, 7.1, 1.5, line_spacing=1.34)
    hairline(s, 8.72, 2.72, 3.75, INK, 1.4)
    text(s, "WHERE THIS COMES FROM", 8.72, 2.86, 3.75, 0.25, size=9.5, font=SANS,
         colour=MIST, bold=True)
    text(s, "Not a prediction. This is what teams running agents in delivery are "
            "reporting right now.",
         8.72, 3.18, 3.75, 1.1, size=12.5, font=SANS, colour=MIST, line_spacing=1.32)

    # 3. What did not move -----------------------------------------------
    s = new(section="The constraint")
    title(s, "The number that did not move is review.")
    col = (CONTENT_W - 0.9) / 2
    for i, (head, items, tint) in enumerate([
        ("Went up, a lot",
         ["Decisions produced per week", "Speed of production", "How polished the output looks"],
         SHELL),
        ("Did not move",
         ["People qualified to judge the decision", "Hours they actually have",
          "Time one honest check takes"],
         SHELL),
    ]):
        left = MARGIN + i * (col + 0.9)
        card(s, left, 2.36, col, 2.28, tint)
        text(s, head.upper(), left + 0.36, 2.62, col - 0.72, 0.3, size=10.5, font=SANS,
             colour=INK if i else MIST, bold=True)
        for j, item in enumerate(items):
            y = 3.06 + j * 0.46
            text(s, item, left + 0.36, y, col - 0.72, 0.4, size=14.5, font=SANS,
                 colour=INK if i else SLATE)
            if j < len(items) - 1:
                hairline(s, left + 0.36, y + 0.36, col - 0.72)
    rich(s, [("When production runs past review capacity, the extra work does not stop. ",
              {"size": 20, "colour": SLATE}),
             ("It gets approved anyway.", {"size": 20, "colour": INK, "bold": True})],
         MARGIN, 5.06, 10.6, 1.0, line_spacing=1.34)

    # 4. Nobody decides to ------------------------------------------------
    s = new(section="The constraint")
    title(s, "Nobody decides to rubber-stamp.")
    text(s, "A reviewer with eleven minutes and nineteen items in the queue produces "
            "nineteen approvals.",
         MARGIN, 2.44, 7.4, 1.0, size=22, font=SERIF, colour=INK, line_spacing=1.3)
    text(s, "Every one is recorded as a review. The governance record is intact. The "
            "governance is not. There is no error message for this, and it looks exactly "
            "like a working PMO right up until something goes wrong and nobody can "
            "reconstruct who checked what.",
         MARGIN, 3.78, 7.4, 1.5, size=15, line_spacing=1.42)
    card(s, 8.72, 2.44, 3.75, 2.5, SHELL)
    text(s, "SILENT DRIFT", 9.06, 2.74, 3.1, 0.3, size=10.5, font=SANS, colour=OVER, bold=True)
    text(s, "The gap between approvals recorded and reviews genuinely done.\n\n"
            "Silent because the system records it as success.",
         9.06, 3.14, 3.1, 1.7, size=14, font=SANS, colour=INK, line_spacing=1.34)
    text(s, "It is not a discipline problem. The reviewer is behaving rationally under a "
            "constraint somebody else handed them.",
         MARGIN, 5.52, 7.4, 0.7, size=15, font=SANS, colour=MIST, italic=True, line_spacing=1.34)

    # 5. The formula ------------------------------------------------------
    s = new(section="The verification budget")
    title(s, "So put a number on it.")
    card(s, MARGIN, 2.34, CONTENT_W, 1.36, SHELL)
    text(s, "VB   =   ( R  x  H  x  u )   /   c", MARGIN, 2.72, CONTENT_W, 0.7,
         size=32, font=SERIF, colour=INK, bold=True, align=PP_ALIGN.CENTER)
    defs = [
        ("R", "People genuinely qualified to judge it. Not everyone with the approve button."),
        ("H", "Hours a week each of them has set aside for reviewing."),
        ("u", "The share of those hours that survives the calendar."),
        ("c", "Hours to genuinely check one decision. The one nobody has."),
    ]
    dw = (CONTENT_W - 0.7) / 2
    for i, (symbol, meaning) in enumerate(defs):
        left = MARGIN + (i % 2) * (dw + 0.7)
        top = 4.14 + (i // 2) * 1.06
        hairline(s, left, top, dw, INK, 1.4)
        text(s, symbol, left, top + 0.14, 0.4, 0.4, size=19, font=SERIF, colour=BRAND, bold=True)
        text(s, meaning, left + 0.52, top + 0.16, dw - 0.52, 0.8, size=13.5, font=SANS,
             colour=SLATE, line_spacing=1.28)
    text(s, "Three of these you can get from a spreadsheet this afternoon.",
         MARGIN, 6.34, CONTENT_W, 0.3, size=13, font=SANS, colour=MIST, italic=True)

    # 6. Six people --------------------------------------------------------
    s = new(section="The verification budget")
    title(s, "Six people. Twenty-one decisions.")
    text(s, "VB   =   ( 6  x  8  x  0.55 )  /  1.25   =   21 a week",
         MARGIN, 2.42, CONTENT_W, 0.5, size=21, font=SERIF, colour=STEEL, bold=True)
    sw = (CONTENT_W - 1.1) / 3
    stat(s, MARGIN, 3.28, sw, "budget", "21",
         "re-baselines a week that can genuinely be reviewed")
    stat(s, MARGIN + sw + 0.55, 3.28, sw, "produced", "70",
         "re-baselines a week the agents actually generate")
    stat(s, MARGIN + (sw + 0.55) * 2, 3.28, sw, "overdraft", "3.3x",
         "about 49 a week carry an approval nobody could defend", OVER)
    hairline(s, MARGIN, 5.62, CONTENT_W)
    text(s, "Six people qualified to review a schedule re-baseline, eight review hours "
            "each, of which 55 percent survives the calendar, and an honest check takes "
            "an hour and a quarter.",
         MARGIN, 5.80, 10.4, 0.8, size=13.5, font=SANS, colour=MIST, line_spacing=1.34)

    # 7. Four classes ------------------------------------------------------
    s = new(section="The four decision classes")
    title(s, "Sort work by how expensive it is to check.")
    # Column left edges and widths. The name column stops short of the test
    # column rather than running under it.
    col_x = [MARGIN + 0.52, 3.62, 8.72, 10.62]
    col_w = [2.10, 4.90, 1.80, 1.85]
    for i, head in enumerate(["CLASS", "THE TEST", "COST TO CHECK", "WHO DECIDES"]):
        text(s, head, col_x[i], 2.46, col_w[i], 0.26, size=9.5, font=SANS,
             colour=MIST, bold=True)
    hairline(s, MARGIN, 2.74, CONTENT_W, INK, 1.4)
    rows = [
        ("A", "machine-checkable", "a test decides it", "minutes", "agent"),
        ("B", "sample-checkable", "check twenty, conclude about two hundred", "minutes, spread", "agent"),
        ("C", "expert-checkable", "a qualified person rebuilds the reasoning", "hours", "agent proposes,\nhuman decides"),
        ("D", "not checkable in advance", "you only find out afterwards", "no answer", "human, always"),
    ]
    y = 2.92
    for letter, name, test, cost, who in rows:
        chip(s, letter, MARGIN, y + 0.03)
        text(s, name, col_x[0], y + 0.05, col_w[0], 0.4, size=13.5, font=SANS,
             colour=INK, bold=True, line_spacing=1.15)
        text(s, test, col_x[1], y + 0.05, col_w[1], 0.4, size=13.5, font=SANS, colour=SLATE)
        text(s, cost, col_x[2], y + 0.05, col_w[2], 0.4, size=13.5, font=SANS, colour=SLATE)
        text(s, who, col_x[3], y + 0.05, col_w[3], 0.62, size=13.5, font=SANS,
             colour=SLATE, line_spacing=1.15)
        y += 0.72
        hairline(s, MARGIN, y - 0.11, CONTENT_W)
    rich(s, [("Not by risk. ", {"bold": True, "colour": INK, "size": 15}),
             ("Risk tells you whether you must check. This tells you what checking costs. "
              "You need both, and most governance frameworks only have the first.",
              {"size": 15, "colour": SLATE})],
         MARGIN, 5.90, 10.9, 0.9, line_spacing=1.34)

    # 8. The quadrant -------------------------------------------------------
    s = new(section="The deployment inversion")
    title(s, "Deploy where checking is cheap,\nnot where the task is easy.", size=30)
    qx, qy, qw, qh = MARGIN, 2.86, 7.4, 3.5
    quad = [
        ("DEPLOY", "Model is good at it and checking is cheap. This is the whole game, and "
                   "it is less glamorous than the other box.", OK, 0, 0),
        ("THE TRAP", "Model is good at it and checking is expensive. Every impressive demo "
                     "lives here. So does every overdraft.", OVER, 1, 0),
        ("WAIT", "Cheap to check, model is not there yet. Harmless. Revisit next year.", MIST, 0, 1),
        ("DO NOT", "Expensive to check and the model is weak. No path from here.", MIST, 1, 1),
    ]
    cw, ch = qw / 2, qh / 2
    for label, blurb, colour, cx, cy in quad:
        left, top = qx + cx * cw, qy + cy * ch
        text(s, label, left + 0.28, top + 0.28, cw - 0.56, 0.3, size=12, font=SANS,
             colour=colour, bold=True)
        text(s, blurb, left + 0.28, top + 0.66, cw - 0.56, 1.1, size=12.5, font=SANS,
             colour=SLATE, line_spacing=1.3)
    hairline(s, qx, qy, qw, INK, 1.4)
    hairline(s, qx, qy + ch, qw)
    hairline(s, qx, qy + qh, qw, INK, 1.4)
    for i in range(3):
        v = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(qx + i * cw), Inches(qy),
                               Emu(int(0.9 * 9525)), Inches(qh))
        v.fill.solid()
        v.fill.fore_color.rgb = INK if i in (0, 2) else LINE
        v.line.fill.background()
        v.shadow.inherit = False
    text(s, "COST OF CHECKING, CHEAP ON THE LEFT", qx, qy - 0.34, qw, 0.26,
         size=9.5, font=SANS, colour=MIST, bold=True)
    text(s, "MODEL CAPABILITY, HIGHER AT THE TOP", qx, qy + qh + 0.16, qw, 0.26,
         size=9.5, font=SANS, colour=MIST, bold=True)

    card(s, 8.6, 2.86, 3.87, 3.5, SHELL)
    text(s, "The top right is where every impressive demo lives.", 8.94, 3.18, 3.2, 0.9,
         size=16, font=SERIF, colour=INK, bold=True, line_spacing=1.26)
    text(s, "Portfolio prioritisation.\nVendor dispute strategy.\nStrategic dependency "
            "renegotiation.\n\nThe models are genuinely good at all three. That is not the "
            "problem. The problem is that your review capacity for them is a couple of "
            "dozen a week.",
         8.94, 4.06, 3.2, 2.2, size=11.5, font=SANS, colour=SLATE, line_spacing=1.30)

    # 10. Part 2 divider ------------------------------------------------------
    s = new(INK, chrome=False)
    part_divider(s, "Part 2", "The operating model",
                 "Processes, roles, governance, technology and metrics, "
                 "redesigned around one constraint.")

    # 11. Five layers -----------------------------------------------------------
    s = new(section="Part 2, overview")
    title(s, "Five layers. Each answers a different question.", size=31)
    table(
        s, 2.44,
        [(MARGIN, 0.4, ""), (MARGIN + 0.45, 2.1, "layer"),
         (3.5, 4.6, "the question it answers"), (8.4, 4.07, "the design rule")],
        [
            ["1", "Processes", "Which decisions go autonomous, and in what order?",
             "Sort by verification cost, not difficulty"],
            ["2", "Roles", "Who owns an agent, and who owns checking it?",
             "An agent role contract, four mandatory fields"],
            ["3", "Governance", "Where does authority stop and escalation start?",
             "Decision rights tied to blast radius"],
            ["4", "Technology", "What must the system emit to make checking cheap?",
             "The evidence plane"],
            ["5", "Metrics", "How do we know supervision is still real?",
             "Six measures, one of them uncomfortable"],
        ],
        pitch=0.62, bold_first=False,
    )
    text(s, "Nothing here is a new department. Four of the five are things you already "
            "have, redesigned around a constraint you did not previously have a number for.",
         MARGIN, 6.02, 11.4, 0.62, size=13, font=SANS, colour=MIST, italic=True,
         line_spacing=1.30)

    # 12. Layer 1, the deployment order ------------------------------------------
    s = new(section="Layer 1, processes")
    title(s, "The deployment order most PMOs use is backwards.", size=31)
    qx, qy, qw, qh = MARGIN, 2.68, 6.9, 2.32
    for label, blurb, colour, cx, cy in [
        ("GO FIRST", "Easy to build, cheap to check. Start here today, and it is the "
                     "only box that scales.", OK, 0, 0),
        ("HOLD", "Easy to build, expensive to check. The trap, and where every "
                 "impressive demo lives.", OVER, 1, 0),
        ("GO SECOND", "Hard to build but cheap to check. Worth the engineering.", MIST, 0, 1),
        ("NEVER YET", "Hard and expensive to check. Keep it human for now.", MIST, 1, 1),
    ]:
        left, top = qx + cx * (qw / 2), qy + cy * (qh / 2)
        text(s, label, left + 0.26, top + 0.22, qw / 2 - 0.52, 0.28, size=11.5,
             font=SANS, colour=colour, bold=True)
        text(s, blurb, left + 0.26, top + 0.56, qw / 2 - 0.52, 0.7, size=12,
             font=SANS, colour=SLATE, line_spacing=1.28)
    hairline(s, qx, qy, qw, INK, 1.4)
    hairline(s, qx, qy + qh / 2, qw)
    hairline(s, qx, qy + qh, qw, INK, 1.4)
    for i in range(3):
        v = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(qx + i * (qw / 2)),
                               Inches(qy), Emu(int(0.9 * 9525)), Inches(qh))
        v.fill.solid()
        v.fill.fore_color.rgb = INK if i in (0, 2) else LINE
        v.line.fill.background()
        v.shadow.inherit = False
    text(s, "TASK DIFFICULTY FOR THE MODEL, EASIER ON THE LEFT", qx, qy - 0.32, qw, 0.26,
         size=9, font=SANS, colour=MIST, bold=True)
    text(s, "VERIFICATION COST, CHEAPER AT THE TOP", qx, qy + qh + 0.14, qw, 0.26,
         size=9, font=SANS, colour=MIST, bold=True)

    card(s, 8.1, 2.68, 4.37, 2.32, SHELL)
    text(s, "WHAT MOST TEAMS DO", 8.42, 2.90, 3.8, 0.26, size=9.5, font=SANS,
         colour=MIST, bold=True)
    text(s, "Deploy along the difficulty axis. Easiest first, hardest last. It feels "
            "rational and it is how every roadmap gets drawn.",
         8.42, 3.18, 3.8, 0.8, size=12, font=SANS, colour=SLATE, line_spacing=1.28)
    text(s, "WHAT ACTUALLY GOVERNS YOU", 8.42, 4.02, 3.8, 0.26, size=9.5, font=SANS,
         colour=MIST, bold=True)
    text(s, "Deploy along the verification axis. You can only run as much autonomy as "
            "you can check.",
         8.42, 4.32, 3.8, 0.7, size=12, font=SANS, colour=SLATE, line_spacing=1.28)
    text(s, "Risk analysis is the easier AI problem. It should go live later than status "
            "reporting.",
         MARGIN, 5.56, 11.4, 0.5, size=17, font=SERIF, colour=INK, bold=True)

    # 13. Four fields --------------------------------------------------------
    s = new(section="The agent role contract")
    title(s, "Every agent gets four fields. Not nine.")
    fields = [
        ("Scope", "Which decisions it may make, named one by one. Anything unnamed is out, "
                  "and running into one is a stop."),
        ("Evidence", "What it must show for every decision. If it cannot produce the "
                     "evidence, it does not make the decision."),
        ("Escalation", "The named conditions where it stops, and the named human it "
                       "stops to."),
        ("Revocation", "How you turn it off, who can, how fast, and what happens to work "
                       "in flight."),
    ]
    fw = (CONTENT_W - 0.6) / 2
    for i, (head, body) in enumerate(fields):
        left = MARGIN + (i % 2) * (fw + 0.6)
        top = 2.46 + (i // 2) * 1.66
        card(s, left, top, fw, 1.36, SHELL)
        text(s, head, left + 0.34, top + 0.24, fw - 0.68, 0.36, size=17, font=SERIF,
             colour=INK, bold=True)
        text(s, body, left + 0.34, top + 0.66, fw - 0.68, 0.66, size=13, font=SANS,
             colour=SLATE, line_spacing=1.3)
    text(s, "Four is the number somebody can hold in their head at six on a Friday when the "
            "thing is misbehaving and a call has to be made. That is the only moment the "
            "contract has to work.",
         MARGIN, 5.96, 11.0, 0.8, size=14.5, font=SANS, colour=MIST, italic=True,
         line_spacing=1.34)

    # 14. Layer 2, what happens to the people ----------------------------------
    s = new(section="Layer 2, roles")
    title(s, "What happens to the people.")
    end_y = table(
        s, 2.46,
        [(MARGIN, 2.7, "role"), (3.75, 3.1, "before"), (7.15, 5.32, "after")],
        [
            ["Project manager", "Coordinator and chaser", "Verifier and exception handler"],
            ["PMO analyst", "Report producer", "Owner of the evidence plane"],
            ["PMO lead", "Process owner", "Owner of decision rights"],
            ["Agent steward", "does not exist",
             "Owns agent contracts, tracks reversals, retires agents"],
            ["Verification lead", "does not exist",
             "Owns the budget, allocates review capacity"],
        ],
        pitch=0.62, size=14, bold_first=True,
    )
    text(s, "Two new roles. Neither is a data scientist. Both are accountability roles, "
            "which is exactly why they get skipped.",
         MARGIN, end_y + 0.24, 11.4, 0.5, size=15, font=SANS, colour=MIST, italic=True)

    # 15. Layer 3, authority by blast radius -------------------------------------
    s = new(section="Layer 3, governance")
    title(s, "Authority is set by blast radius, not by confidence.", size=31)
    end_y = table(
        s, 2.40,
        [(MARGIN, 3.5, "authority level"), (4.55, 4.0, "the test that puts you here"),
         (8.75, 3.72, "examples")],
        [
            ["Agent acts, logs only", "Reversible within a day, nobody outside affected",
             "Reassign a task, update a date"],
            ["Agent acts, human notified", "Reversible within a week, internal only",
             "Reallocate slack, reorder a backlog"],
            ["Agent proposes, one human approves",
             "Reversible with effort, or affects one team plan",
             "Move a milestone, change a resource split"],
            ["Agent proposes, a committee approves",
             "Hard to reverse, or touches cost, contract or compliance",
             "Budget reallocation, vendor change"],
            ["Human decides, agent supports only", "Irreversible, or legally accountable",
             "Contract signature, termination"],
        ],
        pitch=0.58, size=13, bold_first=True,
    )
    rich(s, [("Notice what is absent. ", {"bold": True, "colour": INK, "size": 14}),
             ("Model confidence appears nowhere on this table. Confident and wrong on an "
              "irreversible decision is still a disaster, and a confidence score is the "
              "agent marking its own work.", {"size": 14, "colour": SLATE})],
         MARGIN, end_y + 0.22, 11.4, 0.8, line_spacing=1.30)

    # 16. Six fields ---------------------------------------------------------
    s = new(section="The evidence plane")
    title(s, "Six things every decision has to carry.")
    items = [
        ("1", "What was decided.", "One sentence, naming the thing."),
        ("2", "What it rested on.", "The actual sources, not a summary."),
        ("3", "What was rejected, and why.", "At least one real alternative."),
        ("4", "How confident, and how it would be wrong.", "Named, not hedged."),
        ("5", "How to undo it.", "What that costs, and how long it stays cheap."),
        ("6", "Who is accountable.", "A person, not a team."),
    ]
    iw = (7.9 - 0.6) / 2
    for i, (num, head, body) in enumerate(items):
        left = MARGIN + (i % 2) * (iw + 0.6)
        top = 2.44 + (i // 2) * 1.18
        text(s, num, left, top + 0.02, 0.3, 0.3, size=12, font=SERIF, colour=MIST, bold=True)
        text(s, head, left + 0.34, top, iw - 0.34, 0.58, size=14, font=SANS, colour=INK,
             bold=True, line_spacing=1.18)
        text(s, body, left + 0.34, top + 0.60, iw - 0.34, 0.4, size=12.5, font=SANS,
             colour=MIST, line_spacing=1.26)
        hairline(s, left, top + 1.04, iw)
    card(s, 9.0, 2.44, 3.47, 3.44, INK)
    text(s, "A decision missing any of the six is not a decision.\n\nIt is an output. "
            "Outputs do not get actioned.",
         9.34, 2.86, 2.8, 2.0, size=16, font=SERIF, colour=ON_DARK, line_spacing=1.32)
    text(s, "Every field exists to make the next review faster. The second one alone is "
            "most of the saving.",
         9.34, 4.94, 2.8, 0.9, size=11.5, font=SANS, colour=ON_DARK_MUTED, line_spacing=1.3)

    # 17. Layer 4, reference shape ----------------------------------------------
    s = new(section="Layer 4, technology")
    title(s, "Reference shape, end to end.")
    stages = [
        ("SOURCES", ["Work tracker", "Code and CI", "Finance ledger", "Comms and docs"]),
        ("AGENTS", ["Class A checkers", "Summarisers", "Analysts (propose only)",
                    "Escalation router"]),
        ("EVIDENCE PLANE", ["Decision artifacts", "Source links", "Counter-case",
                            "Reversal path"]),
        ("VERIFICATION", ["Rule engine (auto)", "Sampling queue", "Expert review queue",
                          "Committee record"]),
    ]
    colw = (CONTENT_W - 1.5) / 4
    for i, (head, items) in enumerate(stages):
        left = MARGIN + i * (colw + 0.5)
        text(s, head, left, 2.44, colw, 0.26, size=9.5, font=SANS, colour=MIST, bold=True)
        hairline(s, left, 2.72, colw, INK, 1.4)
        for j, item in enumerate(items):
            text(s, item, left, 2.90 + j * 0.44, colw, 0.36, size=13, font=SANS,
                 colour=SLATE)
        if i < 3:
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                       Inches(left + colw + 0.10), Inches(3.30),
                                       Inches(0.30), Inches(0.20))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = LINE
            arrow.line.fill.background()
            arrow.shadow.inherit = False
    hairline(s, MARGIN, 4.90, CONTENT_W)
    text(s, "The write path back to the source systems runs through verification, never "
            "around it.",
         MARGIN, 5.06, 11.4, 0.4, size=12.5, font=SANS, colour=MIST)
    text(s, "Nothing writes back to the source systems until it has cleared "
            "verification. The write path is the governance boundary.",
         MARGIN, 5.62, 9.6, 0.9, size=19, font=SERIF, colour=INK, line_spacing=1.30)

    # 18. Agentic ------------------------------------------------------------
    s = new(section="Agentic verification")
    title(s, "Agents can supply review capacity,\nnot just consume it.", size=30)
    card(s, MARGIN, 2.86, 5.5, 1.06, SHELL)
    text(s, "c_eff  =  c_a  +  ( 1 - k )  x  c", MARGIN, 3.14, 5.5, 0.5,
         size=22, font=SERIF, colour=INK, bold=True, align=PP_ALIGN.CENTER)
    text(s, "Put a second agent in front of the human. It closes some decisions on its own, "
            "with a machine-checkable reason. The budget goes up without hiring anybody.",
         MARGIN, 4.16, 5.5, 1.2, size=14, font=SANS, colour=SLATE, line_spacing=1.36)
    text(s, "THREE RULES, OR THIS IS WISHFUL THINKING", 6.9, 2.60, 5.6, 0.28,
         size=9.5, font=SANS, colour=MIST, bold=True)
    hairline(s, 6.9, 2.88, 5.57, INK, 1.4)
    rules = [
        ("1", "The verifier output must be machine-checkable.",
         "If a human reads its reasoning, you moved the cost rather than removing it."),
        ("2", "Measure how often it misses a bad decision.",
         "Unmeasured means it counts as zero. Not an estimate. Zero."),
        ("3", "Bank the bottom of the confidence range.",
         "Never the headline number. You are sizing a safety margin."),
    ]
    for i, (num, head, body) in enumerate(rules):
        top = 3.08 + i * 1.14
        text(s, num, 6.9, top + 0.02, 0.34, 0.3, size=13, font=SERIF, colour=MIST, bold=True)
        text(s, head, 7.3, top, 5.17, 0.42, size=14.5, font=SANS, colour=INK, bold=True,
             line_spacing=1.2)
        text(s, body, 7.3, top + 0.40, 5.17, 0.6, size=12.5, font=SANS, colour=MIST,
             line_spacing=1.3)
        if i < 2:
            hairline(s, 6.9, top + 0.98, 5.57)

    # 19. Six numbers ---------------------------------------------------------
    s = new(section="The six operating metrics")
    title(s, "Six numbers. Always per class, never averaged.", size=31)
    metrics = [
        ("Budget", "How much review you actually have", "measure it, do not target it"),
        ("Overdraft", "Whether you are spending capacity you have", "at or under 1.0"),
        ("Silent drift", "Whether the approvals are real", "near baseline"),
        ("Containment", "How much review the agents genuinely supply", "measured, with its miss rate"),
        ("Escalation precision", "Whether escalations are worth the time", "0.7 or better"),
        ("Reversal latency", "Whether the undo button is real", "inside one review cycle"),
    ]
    for i, head in enumerate(["METRIC", "WHAT IT TELLS YOU", "TARGET"]):
        text(s, head, [MARGIN, 4.3, 9.3][i], 2.46, [3.2, 4.8, 3.2][i], 0.26,
             size=9.5, font=SANS, colour=MIST, bold=True)
    hairline(s, MARGIN, 2.74, CONTENT_W, INK, 1.4)
    y = 2.92
    for name, tells, target in metrics:
        highlight = name == "Silent drift"
        text(s, name, MARGIN, y, 3.2, 0.34, size=14.5, font=SANS,
             colour=OVER if highlight else INK, bold=True)
        text(s, tells, 4.3, y, 4.8, 0.34, size=14, font=SANS, colour=SLATE)
        text(s, target, 9.3, y, 3.2, 0.34, size=14, font=SANS, colour=SLATE)
        y += 0.545
        hairline(s, MARGIN, y - 0.10, CONTENT_W)
    text(s, "Average these across a portfolio and a healthy class will hide a drowning one.",
         MARGIN, 6.24, 11.4, 0.42, size=13, font=SANS, colour=MIST, italic=True,
         line_spacing=1.30)

    # 13. Drift ----------------------------------------------------------------
    s = new(section="The six operating metrics")
    title(s, "How to measure the thing nobody measures.", size=31)
    text(s, "Set a floor: the time below which a genuine review could not physically "
            "have happened.",
         MARGIN, 2.40, 7.6, 0.8, size=19, font=SERIF, colour=INK, line_spacing=1.3)
    halves = [
        ("SETTING THE FLOOR",
         ["Watch thirty real reviews. Confirm afterwards that each one was genuine.",
          "Take the tenth percentile of those times.",
          "Never go below reading speed on the document itself."]),
        ("READING THE NUMBER",
         ["A tenth-percentile floor reports about 10 percent even when all is well.",
          "So the signal is the excess above 10, not the number itself.",
          "34 percent means a third of approvals are decorative."]),
    ]
    hw = (7.6 - 0.6) / 2
    for i, (head, lines) in enumerate(halves):
        left = MARGIN + i * (hw + 0.6)
        text(s, head, left, 3.46, hw, 0.26, size=9.5, font=SANS, colour=MIST, bold=True)
        hairline(s, left, 3.74, hw, INK, 1.4)
        for j, item in enumerate(lines):
            text(s, item, left, 3.90 + j * 0.86, hw, 0.8, size=13, font=SANS, colour=SLATE,
                 line_spacing=1.3)
    card(s, 8.9, 2.40, 3.57, 4.0, INK)
    text(s, "Never use this on an individual.", 9.24, 2.78, 2.9, 0.8,
         size=18, font=SERIF, colour=ON_DARK, bold=True, line_spacing=1.24)
    text(s, "The moment you do, review time becomes a thing people manage rather than a "
            "thing you measure. Durations rise, the number goes to zero, and you have lost "
            "the only instrument that can see this failure.\n\n"
            "If your organisation cannot resist, do not collect it at all. A corrupted "
            "metric is worse than a missing one.",
         9.24, 3.70, 2.9, 2.58, size=11, font=SANS, colour=ON_DARK_MUTED, line_spacing=1.32)

    # 21. Layer 5, what the dashboard showed --------------------------------------
    s = new(section="Layer 5, metrics")
    title(s, "Every headline number improved while supervision stopped existing.",
          size=29)
    dw = (CONTENT_W - 1.2) / 4
    for i, (label, value, note) in enumerate([
        ("approval rate", "97%", "up"),
        ("reversal rate", "2%", "down"),
        ("cycle time", "-41%", "down"),
        ("escalations", "steady", "no change"),
    ]):
        stat(s, MARGIN + i * (dw + 0.4), 2.62, dw, label, value, note,
             OK if i < 3 else MIST, 36)
    hairline(s, MARGIN, 4.44, CONTENT_W)
    hw = (CONTENT_W - 0.8) / 2
    for i, (head, body) in enumerate([
        ("WHAT THE DASHBOARD SHOWED",
         "Four green numbers. A programme that looks like it is working, reported to a "
         "board with no reason to doubt it."),
        ("WHAT WAS ACTUALLY HAPPENING",
         "Silent drift climbing every week. The approval rate went up because approving "
         "got faster, and approving got faster because checking stopped."),
    ]):
        left = MARGIN + i * (hw + 0.8)
        text(s, head, left, 4.64, hw, 0.26, size=9.5, font=SANS, colour=MIST, bold=True)
        text(s, body, left, 4.96, hw, 0.9, size=13, font=SANS, colour=SLATE,
             line_spacing=1.30)
    rich(s, [("Put drift on the same page as the good news. ",
              {"bold": True, "colour": INK, "size": 14}),
             ("On its own page nobody looks at it, and the four green numbers win.",
              {"size": 14, "colour": MIST})],
         MARGIN, 6.16, 11.4, 0.5, line_spacing=1.28)

    # 22. Four gates ------------------------------------------------------------
    s = new(section="The four eval gates")
    title(s, "Four gates, each with a number attached.")
    groups = [
        ("BEFORE IT RUNS", [
            ("Do we agree what class this is?",
             "Two qualified people classify fifty decisions separately. If they disagree, "
             "you do not know what checking costs."),
            ("Would it have been right on history?",
             "Replay against decisions you already know the outcome of. Zero Class D taken "
             "on its own."),
        ]),
        ("WHILE IT RUNS", [
            ("Is every decision checkable at the assumed cost?",
             "Continuous, plus twenty read by a human for substance."),
            ("Is the containment real?",
             "Quarterly, against decisions known to be bad."),
        ]),
    ]
    gw = (CONTENT_W - 0.8) / 2
    for i, (head, entries) in enumerate(groups):
        left = MARGIN + i * (gw + 0.8)
        text(s, head, left, 2.44, gw, 0.26, size=9.5, font=SANS, colour=MIST, bold=True)
        hairline(s, left, 2.72, gw, INK, 1.4)
        for j, (q, body) in enumerate(entries):
            top = 2.92 + j * 1.5
            text(s, q, left, top, gw, 0.44, size=15, font=SANS, colour=INK, bold=True,
                 line_spacing=1.24)
            text(s, body, left, top + 0.56, gw, 0.82, size=12.5, font=SANS, colour=SLATE,
                 line_spacing=1.3)
            if j == 0:
                hairline(s, left, top + 1.34, gw)
    text(s, "A gate without a number attached to it is a meeting.",
         MARGIN, 6.10, 11.0, 0.5, size=20, font=SERIF, colour=INK, bold=True)

    # 23. Part 3 divider ----------------------------------------------------------
    s = new(INK, chrome=False)
    part_divider(s, "Part 3", "What to do about it",
                 "Where you are, what is missing, and the smallest useful thing you can "
                 "do on Monday.")

    # 24. Maturity ---------------------------------------------------------------
    s = new(section="Maturity stages")
    title(s, "Five stages, and the test for which one you are on.", size=30)
    stages = [
        ("S0", "unmeasured", "you do not know"),
        ("S1", "measured", "you know the gap"),
        ("S2", "bounded", "you say no to work"),
        ("S3", "contained", "agents supply verified capacity"),
        ("S4", "self-budgeting", "the system throttles itself"),
    ]
    stw = (CONTENT_W - 0.9) / 5
    hairline(s, MARGIN, 2.86, CONTENT_W, INK, 1.4)
    for i, (code, name, blurb) in enumerate(stages):
        left = MARGIN + i * (stw + 0.225)
        text(s, code, left, 2.44, stw, 0.3, size=13, font=SERIF, colour=BRAND, bold=True)
        text(s, name, left, 3.02, stw, 0.4, size=14, font=SANS, colour=INK, bold=True,
             line_spacing=1.18)
        text(s, blurb, left, 3.48, stw, 0.9, size=12, font=SANS, colour=MIST, line_spacing=1.3)
    hairline(s, MARGIN, 4.62, CONTENT_W)
    text(s, "What does it cost you to check one expert decision,\nand when did you last "
            "measure it?",
         MARGIN, 4.86, 7.4, 1.45, size=21, font=SERIF, colour=INK, bold=True,
         line_spacing=1.26)
    text(s, "If there is no number and no date, you are at S0.\n\n"
            "There is no S5. A system with no human verifier does not have an infinite "
            "verification budget. It has none.",
         8.4, 4.86, 4.07, 1.6, size=12.5, font=SANS, colour=MIST, line_spacing=1.32)

    # 16. Limits -----------------------------------------------------------------
    s = new(section="Known limitations")
    title(s, "What I do not know.")
    limits = [
        ("The measurement is biased.",
         "People check more carefully when watched, so the budget always looks better than "
         "it is. Direction known, size not."),
        ("Bursts are not modelled.",
         "The budget is a rate. Work does not arrive evenly. You can be fine on average and "
         "still drop reviews in board week. This is the biggest hole and I am asking for help."),
        ("Timing is a proxy for checking.",
         "Somebody who already knows the answer can approve correctly in twenty seconds."),
        ("The B and C line is genuinely blurry.",
         "Nobody has published how often two qualified people disagree."),
        ("There is no field data.",
         "Every figure here is illustrative. The repository says so on the same page it "
         "prints them."),
    ]
    y = 2.42
    for head, body in limits:
        text(s, head, MARGIN, y, 3.9, 0.5, size=14.5, font=SANS, colour=INK, bold=True,
             line_spacing=1.22)
        text(s, body, 5.0, y, 7.47, 0.6, size=13.5, font=SANS, colour=SLATE, line_spacing=1.3)
        y += 0.86
        hairline(s, MARGIN, y - 0.16, CONTENT_W)

    # 17. Monday ------------------------------------------------------------------
    s = new(section="What to do")
    title(s, "Four things you can do on Monday.")
    steps = [
        ("1", "Pick one decision type an agent already touches.", "One. Not a programme."),
        ("2", "Time thirty reviews of it.",
         "Then ask each reviewer whether they genuinely checked, and discard the ones who "
         "say no. That last step is the whole thing."),
        ("3", "Do the division.",
         "Four numbers into the calculator. About a minute, and it runs in your browser "
         "with nothing sent anywhere."),
        ("4", "Show somebody the gap.", "Not a programme, not a policy. One number, one "
         "conversation."),
    ]
    y = 2.46
    for num, head, body in steps:
        text(s, num, MARGIN, y + 0.02, 0.4, 0.4, size=20, font=SERIF, colour=BRAND, bold=True)
        text(s, head, MARGIN + 0.55, y, 4.5, 0.7, size=15.5, font=SANS, colour=INK, bold=True,
             line_spacing=1.22)
        text(s, body, 5.9, y, 6.57, 0.7, size=13.5, font=SANS, colour=SLATE, line_spacing=1.3)
        y += 0.94
        hairline(s, MARGIN, y - 0.18, CONTENT_W)
    text(s, "Under 1.0 and you have headroom, so deploy more. Over, and you now know "
            "something you did not know this morning. That is the entire first stage.",
         MARGIN, 6.06, 11.4, 0.62, size=13, font=SANS, colour=MIST, italic=True,
         line_spacing=1.30)

    # 18. Closing --------------------------------------------------------------------
    s = new(INK, chrome=False)
    text(s, "Take it and use it", MARGIN, 1.30, 9.0, 1.0,
         size=46, font=SERIF, colour=ON_DARK, bold=True)
    codes = [
        ("qr-repo.png", "Repository", "github.com/hotragn/verb"),
        ("qr-calculator.png", "Calculator", "hotragn.github.io/verb"),
        ("qr-linkedin.png", "LinkedIn", "linkedin.com/in/hotragn-pettugani"),
        ("qr-x.png", "X", "x.com/hotragn"),
    ]
    size = 1.85
    gap = 0.72
    start = MARGIN
    for i, (filename, label, url) in enumerate(codes):
        left = start + i * (size + gap)
        path = ASSETS / filename
        if path.exists():
            s.shapes.add_picture(str(path), Inches(left), Inches(2.72),
                                 Inches(size), Inches(size))
        text(s, label, left, 4.76, size + gap - 0.1, 0.28, size=12.5, font=SANS,
             colour=ON_DARK, bold=True)
        text(s, url, left, 5.04, size + gap - 0.1, 0.28, size=11, font=SANS,
             colour=ON_DARK_ACCENT)
    hairline(s, MARGIN, 5.72, CONTENT_W, DARK_LINE, 1.2)
    text(s, "Formula, schemas, calculator and eval harness. Apache 2.0.\n"
            "Issues and measured verification costs welcome.",
         MARGIN, 5.96, 8.0, 0.8, size=15, font=SERIF, colour=ON_DARK_MUTED, line_spacing=1.34)
    text(s, "Hotragn Pettugani", W - MARGIN - 3.5, 6.16, 3.5, 0.3, size=13, font=SANS,
         colour=ON_DARK, align=PP_ALIGN.RIGHT)

    return prs


def main() -> int:
    if host_logo_path() is None:
        print("  NOTE: no host logo found. Drop the GSDC logo at "
              "deck/assets/gsdc-logo.png and rebuild to place it on every "
              "content slide. The deck is complete without it.")
    missing = [name for name in ("qr-repo.png", "qr-x.png") if not (ASSETS / name).exists()]
    if missing:
        print("QR images missing. Run: python deck/build.py", file=sys.stderr)
        return 1

    prs = build()
    prs.save(OUTPUT)
    print(f"slides   {len(prs.slides.__iter__.__self__._sldIdLst)}")
    print(f"output   {OUTPUT.relative_to(ROOT)}  ({OUTPUT.stat().st_size / 1024:.0f} kB)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
